import fitz
import os
import io
from dotenv import load_dotenv
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

load_dotenv()

# ==== Config ====
pdf_path = os.getenv("DOCUMENT_EDIT_TARGET_PATH")
pptx_output_path = os.getenv("EDITED_DOCUMENT_PATH_PPTX", "output.pptx")
image_dpi = 150  # render DPI
min_font_size = 6
max_font_size = 40
# ================

# --- Step 1: Open PDF & extract text blocks (keep only coords + text) ---
doc = fitz.open(pdf_path)
page = doc[0]
raw_blocks = page.get_text("blocks")
blocks = [(b[0], b[1], b[2], b[3], b[4]) for b in raw_blocks if b[4].strip()]

# --- Step 2: Render PDF page to image ---
pix = page.get_pixmap(dpi=image_dpi)
img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
page_px_w, page_px_h = img.size
page_pt_w, page_pt_h = page.rect.width, page.rect.height
scale = page_px_w / page_pt_w  # pixels per PDF point

# --- Step 3: Auto-detect Y-axis flip ---
def center_px_for_block(block, flip=False):
    x0, y0, x1, y1, _ = block
    cx_pt = (x0 + x1) / 2
    cy_pt = (y0 + y1) / 2
    cx_px = int(round(cx_pt * scale))
    cy_px = int(round((page_pt_h - cy_pt) * scale if flip else cy_pt * scale))
    return max(0, min(page_px_w - 1, cx_px)), max(0, min(page_px_h - 1, cy_px))

def darkness(px, py):
    r, g, b = img.getpixel((px, py))
    return 255*3 - (r + g + b)

flip = False
if blocks:
    nx, ny = center_px_for_block(blocks[0], flip=False)
    fx, fy = center_px_for_block(blocks[0], flip=True)
    flip = darkness(fx, fy) > darkness(nx, ny)

# --- Step 4: Mask original text ---
draw = ImageDraw.Draw(img)
pad_px = max(1, int(round(2 * scale)))
for x0, y0, x1, y1, _ in blocks:
    if flip:
        top = int(round((page_pt_h - y1) * scale))
        left = int(round(x0 * scale))
        bottom = int(round((page_pt_h - y0) * scale))
        right = int(round(x1 * scale))
    else:
        top = int(round(y0 * scale))
        left = int(round(x0 * scale))
        bottom = int(round(y1 * scale))
        right = int(round(x1 * scale))
    draw.rectangle([(left - pad_px, top - pad_px), (right + pad_px, bottom + pad_px)], fill=(255, 255, 255))

bg_path = "page_bg_masked.png"
img.save(bg_path)

# --- Step 5: Create PPTX with same physical size as rendered image ---
page_in_w = page_px_w / image_dpi
page_in_h = page_px_h / image_dpi
prs = Presentation()
prs.slide_width = Inches(page_in_w)
prs.slide_height = Inches(page_in_h)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(bg_path, 0, 0, width=Inches(page_in_w), height=Inches(page_in_h))

# --- Step 6: Add editable text boxes ---
for x0, y0, x1, y1, text in blocks:
    if flip:
        top_px = int(round((page_pt_h - y1) * scale))
        left_px = int(round(x0 * scale))
    else:
        top_px = int(round(y0 * scale))
        left_px = int(round(x0 * scale))
    width_px = int(round((x1 - x0) * scale))
    height_px = int(round((y1 - y0) * scale))

    x_in = left_px / image_dpi
    y_in = top_px / image_dpi
    w_in = width_px / image_dpi
    h_in = height_px / image_dpi

    textbox = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = textbox.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.strip()

    # Dynamic font size
    box_height_pt = h_in * 72
    line_count = text.strip().count("\n") + 1
    scale_factor = 0.45 if line_count == 1 else 0.35
    font_pt = max(min_font_size, min(max_font_size, int(box_height_pt * scale_factor)))
    for run in p.runs:
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor(0, 0, 0)

    # Red border for visibility
    textbox.line.color.rgb = RGBColor(255, 0, 0)
    textbox.line.width = Pt(1)

# --- Step 7: Save PPTX ---
prs.save(pptx_output_path)
print(f"Saved PPTX to: {pptx_output_path}")
print(f"Auto Y-flip used: {flip}")
